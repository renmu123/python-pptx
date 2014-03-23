# encoding: utf-8

"""
Gherkin step implementations for slide-related features.
"""

from __future__ import absolute_import

from behave import given, when, then
from hamcrest import assert_that, equal_to, is_

from pptx import Presentation
from pptx.parts.slide import (
    _SlidePlaceholder, _SlidePlaceholders, _SlideShapeTree
)
from pptx.shapes.picture import Picture
from pptx.shapes.shape import BaseShape

from .helpers import saved_pptx_path, test_pptx


# given ===================================================

@given('a slide having three shapes')
def given_slide_having_three_shapes(context):
    presentation = Presentation(test_pptx('sld-access-shapes'))
    context.slide = presentation.slides[0]


@given('a slide having two placeholders')
def given_slide_having_two_placeholders(context):
    prs = Presentation(test_pptx('sld-access-shapes'))
    context.slide = prs.slides[0]


@given('a slide placeholder collection')
def given_slide_placeholder_collection(context):
    prs = Presentation(test_pptx('sld-access-shapes'))
    context.slide_placeholders = prs.slides[0].placeholders


@given('a slide shape collection')
def given_a_slide_shape_collection(context):
    presentation = Presentation(test_pptx('sld-access-shapes'))
    context.shapes = presentation.slides[0].shapes_new


@given('I have a reference to a blank slide')
def step_given_ref_to_blank_slide(context):
    context.prs = Presentation()
    slide_layout = context.prs.slide_layouts[6]
    context.sld = context.prs.slides.add_slide(slide_layout)


@given('I have a reference to a slide')
def step_given_ref_to_slide(context):
    context.prs = Presentation()
    slide_layout = context.prs.slide_layouts[0]
    context.sld = context.prs.slides.add_slide(slide_layout)


# when ====================================================

@when('I add a new slide')
def step_when_add_slide(context):
    slide_layout = context.prs.slide_masters[0].slide_layouts[0]
    context.prs.slides.add_slide(slide_layout)


# then ====================================================

@then('each slide shape is of the appropriate type')
def then_each_slide_shape_is_of_the_appropriate_type(context):
    shapes = context.shapes
    expected_types = [_SlidePlaceholder, _SlidePlaceholder, Picture]
    for idx, shape in enumerate(shapes):
        assert type(shape) == expected_types[idx], (
            "got \'%s\'" % type(shape).__name__
        )


@then('I can access a shape by index')
def then_can_access_shape_by_index(context):
    shapes = context.shapes
    for idx in range(2):
        shape = shapes[idx]
        assert isinstance(shape, BaseShape)


@then('I can access a slide placeholder by index')
def then_can_access_slide_placeholder_by_index(context):
    slide_placeholders = context.slide_placeholders
    for idx in range(2):
        slide_placeholder = slide_placeholders[idx]
        assert isinstance(slide_placeholder, _SlidePlaceholder)


@then('I can access the placeholder collection of the slide')
def then_can_access_placeholder_collection_of_slide(context):
    slide = context.slide
    slide_placeholders = slide.placeholders
    msg = 'Slide.placeholders not instance of _SlidePlaceholders'
    assert isinstance(slide_placeholders, _SlidePlaceholders), msg


@then('I can access the shape collection of the slide')
def then_can_access_shapes_of_slide(context):
    slide = context.slide
    shapes = slide.shapes_new
    msg = 'Slide.shapes not instance of _SlideShapeTree'
    assert isinstance(shapes, _SlideShapeTree), msg


@then('I can iterate over the shapes')
def then_can_iterate_over_the_shapes(context):
    shapes = context.shapes
    actual_count = 0
    for shape in shapes:
        actual_count += 1
        assert isinstance(shape, BaseShape)
    assert actual_count == 3


@then('I can iterate over the slide placeholders')
def then_can_iterate_over_the_slide_placeholders(context):
    slide_placeholders = context.slide_placeholders
    actual_count = 0
    for slide_placeholder in slide_placeholders:
        actual_count += 1
        assert isinstance(slide_placeholder, _SlidePlaceholder)
    assert actual_count == 2


@then('the length of the shape collection is 3')
def then_len_of_shape_collection_is_3(context):
    slide = context.slide
    shapes = slide.shapes_new
    assert len(shapes) == 3, (
        'expected len(shapes) of 3, got %s' % len(shapes)
    )


@then('the length of the slide placeholder collection is 2')
def then_len_of_placeholder_collection_is_2(context):
    slide = context.slide
    slide_placeholders = slide.placeholders
    assert len(slide_placeholders) == 2, (
        'expected len(slide_placeholders) of 2, got %s' %
        len(slide_placeholders)
    )


@then('the pptx file contains a single slide')
def step_then_pptx_file_contains_single_slide(context):
    prs = Presentation(saved_pptx_path)
    assert_that(len(prs.slides), is_(equal_to(1)))
